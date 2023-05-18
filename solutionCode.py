from docx import Document
from pptx import Presentation
from pptx.util import Inches
import os
import re

def extract_formulas_from_docx(docx_file):
    document = Document(docx_file)

    formulas = []

    for paragraph in document.paragraphs:
        formula_pattern = r"\$.*?\$"
        matches = re.findall(formula_pattern, paragraph.text)

        formulas.extend(matches)

    return formulas

def extract_graphs_from_docx(docx_file, output_directory):
    document = Document(docx_file)

    graphs = []
    graph_counter = 1

    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if run._r.xml.find('<w:drawing>') != -1:
                drawing = run._r
                for item in drawing.getiterator():
                    if 'pic:pic' in item.tag:
                        for child_item in item.getchildren():
                            if 'pic:blipFill' in child_item.tag:
                                graph = child_item.getchildren()[0].attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if graph:
                                    graphs.append(graph)
                                    image_path = os.path.join(output_directory, f"graph_{graph_counter}.png")
                                    document.part.related_parts[graph]._blob.save(image_path)
                                    graph_counter += 1

    return graphs

def create_pptx_presentation(formulas, graphs, output_pptx):
    prs = Presentation()

    for formula in formulas:
        slide_layout = prs.slide_layouts[6]  # Use layout for inserting content with a title and content

        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Formula"

        content_slide = slide.placeholders[1]
        content_slide.text = formula

    for graph in graphs:
        slide_layout = prs.slide_layouts[6]  # Use layout for inserting content with a title and content

        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Graph"

        content_slide = slide.placeholders[1]
        content_slide.text = ""

        image_path = f"graph_{graphs.index(graph) + 1}.png"
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(6), height=Inches(4.5))

    prs.save(output_pptx)

# Usage example
docx_file_path = "Sample.docx"
output_directory = os.getcwd()
output_pptx_path = "output.pptx"

extracted_formulas = extract_formulas_from_docx(docx_file_path)
extracted_graphs = extract_graphs_from_docx(docx_file_path, output_directory)
create_pptx_presentation(extracted_formulas, extracted_graphs, output_pptx_path)

print(f"Presentation saved at {output_pptx_path}.")
